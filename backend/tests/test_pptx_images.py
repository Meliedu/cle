"""Tests for PPTX image captioning path."""

from __future__ import annotations

import io
from unittest.mock import AsyncMock, patch

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def _large_png() -> bytes:
    """Random 128x128 PNG — ~50KB after compression, above the 8KB skip threshold."""
    import os

    img = Image.frombytes("RGB", (128, 128), os.urandom(128 * 128 * 3))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _tiny_png() -> bytes:
    """Solid-color 8x8 PNG — compresses to well under 1KB, below the skip threshold."""
    img = Image.new("RGB", (8, 8), color=(255, 255, 255))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _build_pptx_with_picture(image_blob: bytes) -> bytes:
    """Construct an in-memory PPTX containing one slide with text and a picture."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(slide_layout)

    # Title text
    title = slide.shapes.title
    if title is not None:
        title.text = "Mitosis phases"

    # Picture shape
    slide.shapes.add_picture(
        io.BytesIO(image_blob),
        left=Inches(1),
        top=Inches(2),
        width=Inches(4),
        height=Inches(3),
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.mark.asyncio
async def test_pptx_captions_large_embedded_image(monkeypatch):
    from app.services import parser

    monkeypatch.setattr(parser.settings, "enable_figure_captions", True)
    monkeypatch.setattr(parser.settings, "openrouter_api_key", "test-key")

    blob = _large_png()
    pptx_bytes = _build_pptx_with_picture(blob)

    mock_caption = AsyncMock(return_value="Diagram of mitosis stages")
    with patch("app.services.parser.caption_image", mock_caption):
        result = await parser.parse_document(pptx_bytes, "pptx", "mitosis.pptx")

    mock_caption.assert_awaited()  # captioned at least once
    assert result.pages, "expected at least one slide-worth of content"
    combined = "\n".join(p.text for p in result.pages)
    assert "[Figure: Diagram of mitosis stages]" in combined


@pytest.mark.asyncio
async def test_pptx_skips_tiny_images(monkeypatch):
    """Sub-threshold images (decorative icons) must not trigger VLM calls."""
    from app.services import parser

    monkeypatch.setattr(parser.settings, "enable_figure_captions", True)
    monkeypatch.setattr(parser.settings, "openrouter_api_key", "test-key")

    blob = _tiny_png()  # well below the 8 KB skip threshold
    pptx_bytes = _build_pptx_with_picture(blob)

    mock_caption = AsyncMock(return_value="irrelevant")
    with patch("app.services.parser.caption_image", mock_caption):
        await parser.parse_document(pptx_bytes, "pptx", "tiny.pptx")

    mock_caption.assert_not_awaited()


@pytest.mark.asyncio
async def test_pptx_handles_caption_failure(monkeypatch):
    """A VLM failure must not crash the pipeline — slide text still surfaces."""
    from app.services import parser

    monkeypatch.setattr(parser.settings, "enable_figure_captions", True)
    monkeypatch.setattr(parser.settings, "openrouter_api_key", "test-key")

    blob = _large_png()
    pptx_bytes = _build_pptx_with_picture(blob)

    mock_caption = AsyncMock(return_value=None)  # simulate all failures returning None
    with patch("app.services.parser.caption_image", mock_caption):
        result = await parser.parse_document(pptx_bytes, "pptx", "failed.pptx")

    # Title text should still be present; no figure markers.
    combined = "\n".join(p.text for p in result.pages)
    assert "Mitosis phases" in combined
    assert "[Figure:" not in combined


@pytest.mark.asyncio
async def test_pptx_captions_picture_placeholder_layout(monkeypatch):
    """'Picture with Caption' layout inserts the image as a PlaceholderPicture
    (shape_type == PLACEHOLDER, not PICTURE). The duck-typed .image check must
    still pick it up.
    """
    from app.services import parser

    monkeypatch.setattr(parser.settings, "enable_figure_captions", True)
    monkeypatch.setattr(parser.settings, "openrouter_api_key", "test-key")

    blob = _large_png()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[8])  # Picture with Caption
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Heart anatomy"

    # Find the picture placeholder and insert the image into it.
    picture_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            picture_placeholder = shape
            break
    assert picture_placeholder is not None, "layout 8 should expose a picture placeholder"
    picture_placeholder.insert_picture(io.BytesIO(blob))

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    mock_caption = AsyncMock(return_value="Labeled diagram of a heart")
    with patch("app.services.parser.caption_image", mock_caption):
        result = await parser.parse_document(pptx_bytes, "pptx", "heart.pptx")

    mock_caption.assert_awaited()
    combined = "\n".join(p.text for p in result.pages)
    assert "[Figure: Labeled diagram of a heart]" in combined


@pytest.mark.asyncio
async def test_pptx_dedupes_repeated_images(monkeypatch):
    """Identical images across slides should be captioned once."""
    from app.services import parser

    monkeypatch.setattr(parser.settings, "enable_figure_captions", True)
    monkeypatch.setattr(parser.settings, "openrouter_api_key", "test-key")

    blob = _large_png()

    # Build a deck with two slides, same image on both.
    prs = Presentation()
    layout = prs.slide_layouts[5]
    for i in range(2):
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = f"Slide {i + 1}"
        slide.shapes.add_picture(
            io.BytesIO(blob),
            left=Inches(1),
            top=Inches(2),
            width=Inches(4),
            height=Inches(3),
        )
    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    mock_caption = AsyncMock(return_value="Shared logo")
    with patch("app.services.parser.caption_image", mock_caption):
        await parser.parse_document(pptx_bytes, "pptx", "dup.pptx")

    # Two slides with identical image → one unique VLM call.
    assert mock_caption.await_count == 1
