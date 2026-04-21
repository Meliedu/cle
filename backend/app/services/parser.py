"""Document parser service.

Dispatches to Docling (+VLM figure captions) with PyMuPDF fallback for PDF,
python-docx for DOCX, python-pptx (+VLM image captions) for PPTX, or Whisper
for mp3/mp4 based on file type.
"""

from __future__ import annotations

import asyncio
import hashlib
import io
import logging
import zipfile
from collections import defaultdict
from dataclasses import dataclass, field

import openai

from app.config import settings
from app.services.vlm import caption_image

logger = logging.getLogger(__name__)

DOCX_TYPE = "docx"
PPTX_TYPE = "pptx"
PDF_TYPE = "pdf"
WHISPER_TYPES = {"mp3", "mp4"}

# Hard cap on the total declared-uncompressed size of an Office zip archive.
# python-docx / python-pptx will materialise the entire archive, so a 1 MB
# file claiming 100 GB of uncompressed content can exhaust worker memory.
_MAX_EXPANDED_BYTES = 500 * 1024 * 1024  # 500 MB


def _guard_office_zip(file_data: bytes, filename: str) -> None:
    """Refuse Office archives that claim to expand beyond ``_MAX_EXPANDED_BYTES``.

    Guards DOCX/PPTX uploads against zip-bomb inputs whose declared
    uncompressed size far exceeds the compressed bytes we've accepted at the
    edge. Checked per-entry and cumulatively so a single oversized entry or a
    swarm of small entries both trip the limit.
    """
    try:
        zf = zipfile.ZipFile(io.BytesIO(file_data))
    except zipfile.BadZipFile as exc:
        raise ValueError(f"Office file {filename} is not a valid zip") from exc
    total = 0
    for info in zf.infolist():
        if info.file_size > _MAX_EXPANDED_BYTES:
            raise ValueError(
                f"Office file {filename} contains oversized entry {info.filename}"
            )
        total += info.file_size
        if total > _MAX_EXPANDED_BYTES:
            raise ValueError(f"Office file {filename} expands beyond safe limit")

# Skip tiny embedded assets (template logos, bullet icons) — avoids wasting
# VLM calls on decorative pixels. 4 KB is low enough to keep small chart
# screenshots (matplotlib line plots) while still filtering favicons/bullets.
# Duplicate decorative images are additionally collapsed via sha256 dedup.
_MIN_IMAGE_BYTES = 4 * 1024

# Cap concurrent VLM calls per deck/PDF so Gemini Flash rate limits aren't
# saturated by a single 40-image lecture deck.
_VLM_CONCURRENCY = 4

# Shared prompt so PDF (Docling-invoked) and PPTX (vlm.py-invoked) captions
# drift together. Asks the VLM to transcribe embedded text verbatim before
# describing visuals — so slide exports, screenshots, and scanned pages
# surface their actual content instead of a vague summary.
_FIGURE_PROMPT = (
    "You are extracting content from a university lecture slide, textbook "
    "page, or figure for a student. If the image contains substantial "
    "readable text (slide bullets, screenshot, scanned page), transcribe "
    "ALL visible text VERBATIM first, preserving order and line breaks. "
    "Then, only if there is additional visual content (diagram, chart, "
    "photograph), briefly describe it — labels, axes, data points, "
    "relationships. Omit the description when the image is pure text. "
    "Keep total output under 300 words. Output only the content — no "
    "preamble, no 'This figure shows…'."
)


@dataclass(frozen=True)
class PageContent:
    page_number: int
    text: str


@dataclass(frozen=True)
class ParseResult:
    text: str
    pages: list[PageContent] = field(default_factory=list)
    word_count: int = 0
    page_count: int = 0


async def parse_document(
    file_data: bytes,
    file_type: str,
    filename: str,
) -> ParseResult:
    """Parse a document and return its text content."""
    normalized = file_type.lower().lstrip(".")

    if normalized == PDF_TYPE:
        return await _parse_pdf(file_data, filename)

    if normalized == DOCX_TYPE:
        return await asyncio.wait_for(
            asyncio.to_thread(_parse_docx, file_data, filename),
            timeout=settings.parser_timeout_seconds,
        )

    if normalized == PPTX_TYPE:
        return await asyncio.wait_for(
            _parse_pptx(file_data, filename),
            timeout=settings.parser_timeout_seconds,
        )

    if normalized in WHISPER_TYPES:
        return await asyncio.wait_for(
            _transcribe_with_whisper(file_data, normalized, filename),
            timeout=settings.parser_timeout_seconds,
        )

    raise ValueError(f"Unsupported file type: {file_type}")


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


async def _parse_pdf(file_data: bytes, filename: str) -> ParseResult:
    """Prefer Docling (captures figures via VLM); fall back to pymupdf text-only.

    Each synchronous parser runs under an ``asyncio.wait_for`` wall-clock cap
    (``settings.parser_timeout_seconds``). A stuck Docling/pymupdf call can
    otherwise pin the worker indefinitely — we'd rather fail the job and let
    retry logic reclaim the task than ship a hung worker.

    After primary extraction, a VLM rescue pass fills in pages with very
    little text (typical of scanned PDFs or slide decks exported to raster
    PDF). See :func:`_rescue_low_text_pdf_pages`.
    """
    result: ParseResult | None = None
    if settings.enable_figure_captions and settings.openrouter_api_key:
        try:
            result = await asyncio.wait_for(
                asyncio.to_thread(_parse_pdf_docling, file_data, filename),
                timeout=settings.parser_timeout_seconds,
            )
        except asyncio.TimeoutError:
            logger.warning(
                "Docling parse timed out for %s; falling back to pymupdf", filename
            )
        except Exception:
            logger.exception(
                "Docling parse failed for %s, falling back to pymupdf", filename
            )
    if result is None:
        result = await asyncio.wait_for(
            asyncio.to_thread(_parse_pdf_pymupdf, file_data, filename),
            timeout=settings.parser_timeout_seconds,
        )

    if (
        settings.enable_page_rescue
        and settings.enable_figure_captions
        and settings.openrouter_api_key
    ):
        result = await _rescue_low_text_pdf_pages(result, file_data, filename)
    return result


def _pdf_page_count(file_data: bytes) -> int:
    import pymupdf

    doc = pymupdf.open(stream=file_data, filetype="pdf")
    try:
        return doc.page_count
    finally:
        doc.close()


def _render_pdf_pages(
    file_data: bytes, page_numbers: list[int], dpi: int
) -> list[bytes]:
    """Render specific 1-based PDF page numbers to PNG bytes."""
    import pymupdf

    doc = pymupdf.open(stream=file_data, filetype="pdf")
    try:
        zoom = max(1.0, dpi / 72.0)  # pymupdf native is 72 dpi
        matrix = pymupdf.Matrix(zoom, zoom)
        out: list[bytes] = []
        for pn in page_numbers:
            if pn < 1 or pn > doc.page_count:
                out.append(b"")
                continue
            page = doc.load_page(pn - 1)
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            out.append(pix.tobytes("png"))
        return out
    finally:
        doc.close()


async def _rescue_low_text_pdf_pages(
    result: ParseResult, file_data: bytes, filename: str
) -> ParseResult:
    """Replace near-empty pages with a verbatim VLM transcription.

    A page is a rescue candidate when the primary parser extracted fewer than
    :attr:`Settings.page_rescue_min_words` words for it, which is the pattern
    for scanned textbook pages or slide decks exported as raster-only PDFs.
    Rescues are capped at :attr:`Settings.page_rescue_max_pages` per document
    to bound VLM cost on long image-heavy files; pages beyond the cap keep
    their original (empty) text.
    """
    # Lazy import to avoid a parser↔vlm import cycle at module load.
    from app.services.vlm import transcribe_page

    try:
        doc_page_count = await asyncio.to_thread(_pdf_page_count, file_data)
    except Exception:
        logger.exception("Page rescue: could not read page count for %s", filename)
        return result

    page_text: dict[int, str] = {p.page_number: p.text for p in result.pages}
    min_words = settings.page_rescue_min_words
    targets = [
        pn
        for pn in range(1, doc_page_count + 1)
        if len(page_text.get(pn, "").split()) < min_words
    ]
    if not targets:
        return result
    capped = targets[: settings.page_rescue_max_pages]

    try:
        rendered = await asyncio.to_thread(
            _render_pdf_pages, file_data, capped, settings.page_rescue_render_dpi
        )
    except Exception:
        logger.exception("Page rescue: render failed for %s", filename)
        return result

    # Bound concurrent VLM calls so a 40-page scanned PDF doesn't fire 40
    # simultaneous requests at OpenRouter. Matches _VLM_CONCURRENCY used by
    # docling's figure-caption path.
    sem = asyncio.Semaphore(_VLM_CONCURRENCY)

    async def _bounded(img: bytes) -> str | None:
        if not img:
            return None
        async with sem:
            return await transcribe_page(img)

    transcriptions = await asyncio.gather(*(_bounded(img) for img in rendered))

    rescued = 0
    for pn, transcribed in zip(capped, transcriptions):
        if transcribed:
            page_text[pn] = transcribed
            rescued += 1

    logger.info(
        "Page rescue: %d/%d low-text pages transcribed (total=%d, filename=%s)",
        rescued,
        len(capped),
        doc_page_count,
        filename,
    )

    new_pages = [
        PageContent(page_number=pn, text=page_text[pn])
        for pn in sorted(page_text)
        if page_text[pn]
    ]
    full_text = "\n\n".join(p.text for p in new_pages)
    return ParseResult(
        text=full_text,
        pages=new_pages,
        word_count=len(full_text.split()),
        page_count=doc_page_count or result.page_count,
    )




def _parse_pdf_pymupdf(file_data: bytes, filename: str) -> ParseResult:
    """Extract plain text per page via raw pymupdf.

    Kept as fallback: fast, no model dependencies, drops images silently.
    """
    import pymupdf

    doc = pymupdf.open(stream=file_data, filetype="pdf")
    try:
        pages: list[PageContent] = []
        full_parts: list[str] = []

        for page_index in range(doc.page_count):
            page = doc.load_page(page_index)
            text = page.get_text("text").strip()
            if text:
                pages.append(PageContent(page_number=page_index + 1, text=text))
                full_parts.append(text)

        full_text = "\n\n".join(full_parts)
        return ParseResult(
            text=full_text,
            pages=pages,
            word_count=len(full_text.split()),
            page_count=len(pages) or doc.page_count or 1,
        )
    finally:
        doc.close()


_DOCLING_CONVERTER = None


def _get_docling_converter():
    """Lazily build + cache a shared Docling ``DocumentConverter``.

    Building a converter loads plugin registries, OCR/layout/tablestructure
    engines, and pipeline options — roughly 3-5 s of wall clock on CPU.
    Worker processes tasks one-at-a-time (see ``app.services.worker``), so
    a single long-lived instance per worker is safe without extra locking
    and saves that init cost on every PDF after the first.
    """
    global _DOCLING_CONVERTER  # noqa: PLW0603
    if _DOCLING_CONVERTER is not None:
        return _DOCLING_CONVERTER

    # Lazy imports: keeps torch / transformers / docling off the import path
    # of services (API, tests) that never parse PDFs.
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import (
        PdfPipelineOptions,
        PictureDescriptionApiOptions,
    )
    from docling.document_converter import DocumentConverter, PdfFormatOption

    api_url = f"{settings.openrouter_base_url.rstrip('/')}/chat/completions"
    picture_opts = PictureDescriptionApiOptions(
        url=api_url,
        headers={"Authorization": f"Bearer {settings.openrouter_api_key}"},
        params={"model": settings.vlm_model, "max_tokens": 600},
        prompt=_FIGURE_PROMPT,
        timeout=float(settings.vlm_timeout_seconds),
        concurrency=_VLM_CONCURRENCY,
    )

    # generate_picture_images left at default (False): we only consume
    # `picture.annotations[*].text`; keeping the rendered image on every
    # PictureItem inflates RAM on figure-heavy PDFs with no upside.
    # do_ocr=False: course PDFs are text-native (lectures, papers); we get
    # figure content via VLM captions, not OCR of embedded text images. Leaving
    # OCR on makes docling load rapidocr, which tries to download ONNX weights
    # into site-packages at runtime — the non-root app user can't write there
    # and the whole pipeline fails.
    pipeline_options = PdfPipelineOptions(
        do_ocr=False,
        do_picture_description=True,
        picture_description_options=picture_opts,
        images_scale=2.0,
        enable_remote_services=True,
    )

    _DOCLING_CONVERTER = DocumentConverter(
        format_options={
            InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options),
        }
    )
    return _DOCLING_CONVERTER


def _parse_pdf_docling(file_data: bytes, filename: str) -> ParseResult:
    """Use Docling for layout-aware text + VLM figure captions.

    Captions are inlined as ``[Figure: ...]`` inside the page's text so the
    chunker keeps them adjacent to the paragraph that introduces them. This
    is intentionally a sync function — it's dispatched via ``to_thread`` so
    Docling's blocking work (torch inference, HTTP to OpenRouter) doesn't
    stall the event loop. Runs only inside the async worker, not the request
    path, so 2–10 s/page latency is acceptable.
    """
    from docling_core.types.io import DocumentStream

    converter = _get_docling_converter()

    stream = DocumentStream(name=filename, stream=io.BytesIO(file_data))
    result = converter.convert(stream)
    document = result.document

    # Group captions by page number so they appear inline on the right page.
    captions_by_page: dict[int, list[str]] = defaultdict(list)
    pictures = getattr(document, "pictures", None) or []
    for picture in pictures:
        caption = _extract_picture_caption(picture)
        if not caption:
            continue
        page_no = _picture_page_no(picture)
        if page_no is not None:
            captions_by_page[page_no].append(caption)

    # Build per-page text. Docling's export_to_markdown is per-document; we want
    # per-page to preserve page_number attribution downstream.
    pages: list[PageContent] = []
    page_objs = getattr(document, "pages", None) or {}

    # document.pages is a dict keyed by page_no; iterate sorted.
    page_items = (
        sorted(page_objs.items()) if isinstance(page_objs, dict) else list(enumerate(page_objs, 1))
    )

    if page_items:
        for page_no, _page in page_items:
            text_parts: list[str] = []
            page_text = _page_text(document, page_no)
            if page_text:
                text_parts.append(page_text)
            for caption in captions_by_page.get(page_no, []):
                text_parts.append(f"[Figure: {caption}]")
            combined = "\n\n".join(text_parts).strip()
            if combined:
                pages.append(PageContent(page_number=int(page_no), text=combined))
    else:
        # Fallback: flatten to one page using markdown export.
        md = document.export_to_markdown().strip()
        if md:
            pages.append(PageContent(page_number=1, text=md))

    full_text = "\n\n".join(p.text for p in pages)
    total_pictures = len(pictures)
    total_captions = sum(len(v) for v in captions_by_page.values())
    logger.info(
        "docling parsed %s: %d pages, %d figures (%d captioned)",
        filename,
        len(pages),
        total_pictures,
        total_captions,
    )
    return ParseResult(
        text=full_text,
        pages=pages,
        word_count=len(full_text.split()),
        page_count=len(pages) or 1,
    )


def _extract_picture_caption(picture) -> str | None:
    """Pull the VLM-generated description out of a Docling PictureItem."""
    annotations = getattr(picture, "annotations", None) or []
    for ann in annotations:
        text = getattr(ann, "text", None)
        if text:
            return str(text).strip()
    return None


def _picture_page_no(picture) -> int | None:
    prov = getattr(picture, "prov", None) or []
    if not prov:
        return None
    page_no = getattr(prov[0], "page_no", None)
    return int(page_no) if page_no is not None else None


def _page_text(document, page_no: int) -> str:
    """Collect text items whose first provenance is on ``page_no``."""
    texts = getattr(document, "texts", None) or []
    out: list[str] = []
    for item in texts:
        prov = getattr(item, "prov", None) or []
        if not prov:
            continue
        if getattr(prov[0], "page_no", None) != page_no:
            continue
        text = getattr(item, "text", None)
        if text:
            out.append(str(text).strip())
    return "\n".join(s for s in out if s)


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def _parse_docx(file_data: bytes, filename: str) -> ParseResult:
    """Extract text from DOCX using python-docx."""
    _guard_office_zip(file_data, filename)
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(file_data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    full_text = "\n\n".join(paragraphs)

    return ParseResult(
        text=full_text,
        pages=[PageContent(page_number=1, text=full_text)],
        word_count=len(full_text.split()),
        page_count=1,
    )


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


async def _parse_pptx(file_data: bytes, filename: str) -> ParseResult:
    """Extract slide text and — where enabled — caption embedded Picture shapes."""
    text_by_slide, image_jobs = await asyncio.to_thread(
        _collect_pptx_content, file_data, filename
    )

    captions_by_slide: dict[int, list[str]] = defaultdict(list)

    if image_jobs and settings.enable_figure_captions and settings.openrouter_api_key:
        # Dedup identical images (template logos repeated on every slide).
        unique: dict[str, tuple[int, bytes, str]] = {}
        per_slide_hashes: list[tuple[int, str]] = []
        for slide_num, blob, context in image_jobs:
            h = hashlib.sha256(blob).hexdigest()
            per_slide_hashes.append((slide_num, h))
            if h not in unique:
                unique[h] = (slide_num, blob, context)

        keys = list(unique.keys())
        # Bound concurrent VLM calls so a 40-image deck doesn't saturate
        # Gemini Flash's per-minute rate limit and shed captions to 429s.
        semaphore = asyncio.Semaphore(_VLM_CONCURRENCY)

        async def _capped(blob: bytes, context: str) -> str | None:
            async with semaphore:
                return await caption_image(blob, context=context)

        coros = [_capped(unique[k][1], unique[k][2]) for k in keys]
        results = await asyncio.gather(*coros, return_exceptions=True)

        caption_by_hash: dict[str, str] = {}
        ok = 0
        for key, res in zip(keys, results):
            if isinstance(res, str) and res:
                caption_by_hash[key] = res
                ok += 1

        for slide_num, h in per_slide_hashes:
            cap = caption_by_hash.get(h)
            if cap:
                captions_by_slide[slide_num].append(cap)

        logger.info(
            "pptx %s: captioned %d/%d images (%d unique)",
            filename,
            ok,
            len(image_jobs),
            len(unique),
        )

    pages: list[PageContent] = []
    for slide_num in sorted(set(text_by_slide) | set(captions_by_slide)):
        parts: list[str] = []
        body = text_by_slide.get(slide_num)
        if body:
            parts.append(body)
        for cap in captions_by_slide.get(slide_num, []):
            parts.append(f"[Figure: {cap}]")
        combined = "\n\n".join(parts).strip()
        if combined:
            pages.append(PageContent(page_number=slide_num, text=combined))

    full_text = "\n\n".join(p.text for p in pages)
    return ParseResult(
        text=full_text,
        pages=pages,
        word_count=len(full_text.split()),
        page_count=len(pages) or 1,
    )


def _iter_picture_blobs(shapes):
    """Yield raw image bytes for every picture-like shape, recursively.

    Duck-types on ``shape.image.blob`` rather than filtering by
    ``MSO_SHAPE_TYPE.PICTURE`` — that misses pictures embedded in placeholder
    layouts ("Picture with Caption"), which are PlaceholderPicture instances
    with ``shape_type == PLACEHOLDER``. Also recurses into GROUP shapes so
    grouped diagrams don't get dropped.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_picture_blobs(shape.shapes)
            continue
        image = getattr(shape, "image", None)
        if image is None:
            continue
        try:
            blob = image.blob
        except Exception:  # noqa: BLE001 — linked images raise here
            continue
        if blob:
            yield blob


def _collect_pptx_content(
    file_data: bytes,
    filename: str = "<pptx>",
) -> tuple[dict[int, str], list[tuple[int, bytes, str]]]:
    """Walk slides synchronously, return (text_by_slide, image_jobs).

    image_jobs: list of (slide_num, image_bytes, surrounding_slide_text).
    """
    _guard_office_zip(file_data, filename)
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_data))
    text_by_slide: dict[int, str] = {}
    image_jobs: list[tuple[int, bytes, str]] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

        slide_text = "\n".join(texts)
        if slide_text:
            text_by_slide[slide_num] = slide_text

        for blob in _iter_picture_blobs(slide.shapes):
            if len(blob) < _MIN_IMAGE_BYTES:
                continue
            image_jobs.append((slide_num, blob, slide_text))

    return text_by_slide, image_jobs


# ---------------------------------------------------------------------------
# Whisper
# ---------------------------------------------------------------------------


async def _transcribe_with_whisper(
    file_data: bytes,
    file_type: str,
    filename: str,
) -> ParseResult:
    """Transcribe mp3/mp4 using OpenAI Whisper."""
    client = openai.AsyncOpenAI(api_key=settings.openai_api_key)

    audio_file = io.BytesIO(file_data)
    audio_file.name = filename

    transcript = await client.audio.transcriptions.create(
        model="whisper-1",
        file=audio_file,
        response_format="text",
    )

    text = str(transcript)
    word_count = len(text.split())

    return ParseResult(
        text=text,
        pages=[PageContent(page_number=1, text=text)],
        word_count=word_count,
        page_count=1,
    )
